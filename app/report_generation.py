"""
템플릿 불러오기, 페이지 파악, 태그 감지 및 처리
"""
import re
import json
from pathlib import Path
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.text_generation import generate_text_for_tag
from app.visualization import create_chart_for_tag


def load_tag_config(config_path: str = 'app/tag_config.json') -> Dict[str, Any]:
    """
    태그 설정 파일 로드
    
    Args:
        config_path: 설정 파일 경로
        
    Returns:
        태그별 설정 딕셔너리
    """
    config_file = Path(config_path)
    if config_file.exists():
        with open(config_file, 'r', encoding='utf-8') as f:
            return json.load(f)
    else:
        print(f"⚠️ 설정 파일이 없습니다: {config_path}")
        return {}


def find_tags_in_slide(slide) -> List[Dict[str, Any]]:
    """
    슬라이드에서 모든 태그를 찾아서 반환
    
    Args:
        slide: pptx 슬라이드 객체
        
    Returns:
        태그 정보 리스트 [{'tag': 'TITLE_AREA', 'shape': shape, 'type': 'text'}, ...]
    """
    tags = []
    tag_pattern = re.compile(r'\{\{(\w+)\}\}')  # {{TAG_NAME}} 패턴
    
    for shape in slide.shapes:
        # shape 이름에서 태그 찾기
        if hasattr(shape, 'name') and shape.name:
            matches = tag_pattern.findall(shape.name)
            for tag in matches:
                tags.append({
                    'tag': tag,
                    'shape': shape,
                    'type': 'name',
                    'original_text': shape.name
                })
        
        # 텍스트 프레임에서 태그 찾기
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text = shape.text_frame.text
            matches = tag_pattern.findall(text)
            for tag in matches:
                tags.append({
                    'tag': tag,
                    'shape': shape,
                    'type': 'text',
                    'original_text': text
                })
    
    return tags


def apply_text_styling(shape, config: Dict[str, Any]) -> None:
    """
    shape에 텍스트 스타일 적용
    
    Args:
        shape: pptx shape 객체
        config: 태그 설정 딕셔너리
    """
    if not hasattr(shape, 'text_frame'):
        return
    
    # 모든 문단과 런에 스타일 적용
    for paragraph in shape.text_frame.paragraphs:
        # 정렬
        if 'alignment' in config:
            align_map = {
                'left': PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right': PP_ALIGN.RIGHT,
                'justify': PP_ALIGN.JUSTIFY
            }
            if config['alignment'].lower() in align_map:
                paragraph.alignment = align_map[config['alignment'].lower()]
        
        # 각 런(run)에 폰트 스타일 적용
        for run in paragraph.runs:
            # 폰트 크기
            if 'font_size' in config:
                run.font.size = Pt(config['font_size'])
            
            # 볼드
            if 'font_bold' in config:
                run.font.bold = config['font_bold']
            
            # 폰트 색상
            if 'font_color' in config:
                color = config['font_color']
                if isinstance(color, list) and len(color) == 3:
                    run.font.color.rgb = RGBColor(color[0], color[1], color[2])
        
        # 런이 없는 경우(빈 문단)에도 기본 폰트 설정
        if len(paragraph.runs) == 0 and paragraph.text:
            run = paragraph.add_run()
            run.text = paragraph.text
            paragraph.text = ""  # 기존 텍스트 제거
            
            # 폰트 크기
            if 'font_size' in config:
                run.font.size = Pt(config['font_size'])
            
            # 볼드
            if 'font_bold' in config:
                run.font.bold = config['font_bold']
            
            # 폰트 색상
            if 'font_color' in config:
                color = config['font_color']
                if isinstance(color, list) and len(color) == 3:
                    run.font.color.rgb = RGBColor(color[0], color[1], color[2])


def process_tag(
    tag_info: Dict[str, Any],
    tag_config: Dict[str, Any],
    phrase_data: List[Dict[str, Any]],
    month: str,
    gemini_api_key: str,
    slide,
    context: Optional[Dict[str, Any]] = None
) -> None:
    """
    개별 태그 처리
    
    Args:
        tag_info: 태그 정보 {'tag': 'TITLE_AREA', 'shape': shape, ...}
        tag_config: 전체 태그 설정
        phrase_data: 프레이즈 데이터
        month: 월 정보
        gemini_api_key: Gemini API 키
        slide: 슬라이드 객체
    """
    tag_name = tag_info['tag']
    shape = tag_info['shape']
    
    # 스킵할 태그들 (GEMINI 호출 없이 스킵)
    skip_tags = ['ANALYSIS_AREA', 'PRODUCT_AREA']
    if tag_name in skip_tags:
        print(f"  ⏭️ 태그 스킵: {tag_name}")
        return
    
    # 태그별 설정 가져오기
    config = tag_config.get(tag_name, {})
    
    # 태그 타입에 따라 처리
    tag_type = config.get('type', 'text')  # 'text', 'chart', 'list', 'asset', 'composite'
    
    if tag_type == 'chart':
        # 차트를 슬라이드에 직접 삽입 (원으로 그리기)
        create_chart_for_tag(
            tag_name=tag_name,
            phrase_data=phrase_data,
            config=config,
            slide=slide,
            marker_text=f"{{{{{tag_name}}}}}"
        )
        print(f"📊 차트 삽입 완료: {tag_name}")
    elif tag_type == 'asset':
        # 에셋 생성 (워드클라우드, 이미지 등)
        # TODO: 에셋 생성 로직
        print(f"🖼️ 에셋 생성 필요: {tag_name}")
    elif tag_type == 'composite':
        # 복합 타입 처리
        # TODO: composite 타입 처리 로직
        print(f"📋 복합 타입 처리 필요: {tag_name}")
    elif tag_type == 'list':
        # 리스트 타입 처리
        text = generate_text_for_tag(
            tag_name=tag_name,
            phrase_data=phrase_data,
            month=month,
            gemini_api_key=gemini_api_key,
            config=config,
            context=context
        )
        # 리스트 형식으로 변환 필요할 수 있음
        if hasattr(shape, 'text_frame'):
            # 기존 텍스트 프레임 초기화 후 새 텍스트 추가
            text_frame = shape.text_frame
            text_frame.clear()
            
            # 첫 번째 문단에 텍스트 추가
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = text
            
            apply_text_styling(shape, config)
            print(f"✅ 리스트 삽입: {tag_name}")
    else:
        # 텍스트 생성 (기본)
        # TITLE_AREA와 SUBTITLE1_AREA는 prompt_template을 직접 사용 (AI 생성 없음)
        if tag_name in ['TITLE_AREA', 'SUBTITLE1_AREA']:
            # prompt_template에서 변수만 치환해서 직접 사용
            prompt_template = config.get('prompt_template', '')
            format_vars = {'month': month}
            try:
                text = prompt_template.format(**format_vars)
            except KeyError:
                # 변수 치환 실패 시 원본 사용
                text = prompt_template
        else:
            # 다른 태그는 AI로 텍스트 생성
            text = generate_text_for_tag(
                tag_name=tag_name,
                phrase_data=phrase_data,
                month=month,
                gemini_api_key=gemini_api_key,
                config=config,
                context=context
            )
        
        # shape에 텍스트 삽입
        if hasattr(shape, 'text_frame'):
            # TITLE_AREA는 기존 포맷팅 유지
            if tag_name == 'TITLE_AREA':
                # 기존 포맷팅을 유지하면서 텍스트만 교체
                text_frame = shape.text_frame
                if len(text_frame.paragraphs) > 0:
                    # 첫 번째 문단의 첫 번째 런에 텍스트만 교체
                    paragraph = text_frame.paragraphs[0]
                    if len(paragraph.runs) > 0:
                        # 기존 런의 포맷팅 유지하면서 텍스트만 변경
                        paragraph.runs[0].text = text
                        # 추가 런이 있으면 역순으로 제거 (인덱스 변경 방지)
                        runs_to_remove = list(paragraph.runs[1:])
                        for run in reversed(runs_to_remove):
                            paragraph._element.remove(run._element)
                    else:
                        # 런이 없으면 기존 포맷팅 복사해서 새 런 추가
                        run = paragraph.add_run()
                        run.text = text
                else:
                    # 문단이 없으면 새로 추가 (기본 포맷팅 사용)
                    paragraph = text_frame.paragraphs[0]
                    run = paragraph.add_run()
                    run.text = text
            elif tag_name == 'SUBTITLE1_AREA':
                # SUBTITLE1_AREA는 텍스트 교체 후 config의 폰트 사이즈 적용
                text_frame = shape.text_frame
                if len(text_frame.paragraphs) > 0:
                    paragraph = text_frame.paragraphs[0]
                    if len(paragraph.runs) > 0:
                        # 기존 런의 텍스트만 변경
                        paragraph.runs[0].text = text
                        # 추가 런이 있으면 역순으로 제거
                        runs_to_remove = list(paragraph.runs[1:])
                        for run in reversed(runs_to_remove):
                            paragraph._element.remove(run._element)
                    else:
                        run = paragraph.add_run()
                        run.text = text
                else:
                    paragraph = text_frame.paragraphs[0]
                    run = paragraph.add_run()
                    run.text = text
                
                # config의 폰트 사이즈 적용
                apply_text_styling(shape, config)
            elif tag_name == 'DESCRIPTION2_AREA':
                # DESCRIPTION2_AREA는 텍스트 교체 후 config의 폰트 사이즈 적용
                text_frame = shape.text_frame
                if len(text_frame.paragraphs) > 0:
                    paragraph = text_frame.paragraphs[0]
                    if len(paragraph.runs) > 0:
                        # 기존 런의 텍스트만 변경
                        paragraph.runs[0].text = text
                        # 추가 런이 있으면 역순으로 제거
                        runs_to_remove = list(paragraph.runs[1:])
                        for run in reversed(runs_to_remove):
                            paragraph._element.remove(run._element)
                    else:
                        run = paragraph.add_run()
                        run.text = text
                else:
                    paragraph = text_frame.paragraphs[0]
                    run = paragraph.add_run()
                    run.text = text
                
                # config의 폰트 사이즈 적용
                apply_text_styling(shape, config)
            else:
                # 다른 태그는 기존 방식대로 포맷팅 적용
                text_frame = shape.text_frame
                text_frame.clear()
                
                # 첫 번째 문단에 텍스트 추가
                paragraph = text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = text
                
                # 스타일 적용
                apply_text_styling(shape, config)
            
            print(f"✅ 텍스트 삽입: {tag_name}")


def generate_report(
    template_path: str,
    output_path: str,
    phrase_data: List[Dict[str, Any]],
    month: str,
    gemini_api_key: str,
    tag_config_path: str = 'app/tag_config.json',
    weather_analysis: Optional[Dict[str, Any]] = None
) -> None:
    """
    리포트 생성 메인 함수
    
    Args:
        template_path: 템플릿 PPT 파일 경로
        output_path: 출력 PPT 파일 경로
        phrase_data: 프레이즈 데이터
        month: 월 정보
        gemini_api_key: Gemini API 키
        tag_config_path: 태그 설정 파일 경로
        weather_analysis: 기상 데이터 분석 결과
    """
    print(f"📄 템플릿 로드: {template_path}")
    prs = Presentation(template_path)
    
    # 태그 설정 로드
    tag_config = load_tag_config(tag_config_path)
    
    # 출력 디렉토리 생성
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    
    print(f"📑 총 슬라이드 수: {len(prs.slides)}")
    
    # 각 슬라이드 처리
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 슬라이드 {slide_idx + 1} 처리 중...")
        
        # 슬라이드에서 모든 태그 찾기
        tags = find_tags_in_slide(slide)
        
        if not tags:
            print(f"  ℹ️ 태그 없음")
            continue
        
        print(f"  🔍 발견된 태그: {[t['tag'] for t in tags]}")
        
        # 슬라이드 단위 context (태그 간 데이터 공유용)
        slide_context = {}
        
        # 기상 데이터를 context에 추가
        if weather_analysis:
            slide_context['weather_analysis'] = weather_analysis
        
        # KEYWORD1_AREA를 먼저 처리하여 insight_title 생성
        keyword1_tag = next((t for t in tags if t['tag'] == 'KEYWORD1_AREA'), None)
        keyword1_text = None
        if keyword1_tag:
            try:
                keyword1_config = tag_config.get('KEYWORD1_AREA', {})
                keyword1_text = generate_text_for_tag(
                    tag_name='KEYWORD1_AREA',
                    phrase_data=phrase_data,
                    month=month,
                    gemini_api_key=gemini_api_key,
                    config=keyword1_config,
                    context=slide_context
                )
                # KEYWORD1_AREA 결과를 insight_title로 저장
                slide_context['insight_title'] = keyword1_text
                print(f"  📌 KEYWORD1_AREA 결과를 insight_title로 저장: {keyword1_text}")
            except Exception as e:
                print(f"  ⚠️ KEYWORD1_AREA 처리 실패: {e}")
        
        # KEYWORD2_AREA를 먼저 처리하여 insight_title2 생성
        keyword2_tag = next((t for t in tags if t['tag'] == 'KEYWORD2_AREA'), None)
        keyword2_text = None
        if keyword2_tag:
            try:
                keyword2_config = tag_config.get('KEYWORD2_AREA', {})
                keyword2_text = generate_text_for_tag(
                    tag_name='KEYWORD2_AREA',
                    phrase_data=phrase_data,
                    month=month,
                    gemini_api_key=gemini_api_key,
                    config=keyword2_config,
                    context=slide_context
                )
                # KEYWORD2_AREA 결과를 insight_title2로 저장
                slide_context['insight_title2'] = keyword2_text
                print(f"  📌 KEYWORD2_AREA 결과를 insight_title2로 저장: {keyword2_text}")
            except Exception as e:
                print(f"  ⚠️ KEYWORD2_AREA 처리 실패: {e}")
        
        # 각 태그 처리
        for tag_info in tags:
            # KEYWORD1_AREA는 이미 처리했으므로 텍스트만 삽입
            if tag_info['tag'] == 'KEYWORD1_AREA' and keyword1_text:
                try:
                    shape = tag_info['shape']
                    config = tag_config.get('KEYWORD1_AREA', {})
                    if hasattr(shape, 'text_frame'):
                        text_frame = shape.text_frame
                        text_frame.clear()
                        paragraph = text_frame.paragraphs[0]
                        run = paragraph.add_run()
                        run.text = keyword1_text
                        apply_text_styling(shape, config)
                        print(f"✅ 텍스트 삽입: KEYWORD1_AREA")
                except Exception as e:
                    print(f"  ❌ KEYWORD1_AREA 삽입 실패: {e}")
            # KEYWORD2_AREA는 이미 처리했으므로 텍스트만 삽입
            elif tag_info['tag'] == 'KEYWORD2_AREA' and keyword2_text:
                try:
                    shape = tag_info['shape']
                    config = tag_config.get('KEYWORD2_AREA', {})
                    if hasattr(shape, 'text_frame'):
                        text_frame = shape.text_frame
                        text_frame.clear()
                        paragraph = text_frame.paragraphs[0]
                        run = paragraph.add_run()
                        run.text = keyword2_text
                        apply_text_styling(shape, config)
                        print(f"✅ 텍스트 삽입: KEYWORD2_AREA")
                except Exception as e:
                    print(f"  ❌ KEYWORD2_AREA 삽입 실패: {e}")
            else:
                try:
                    process_tag(
                        tag_info=tag_info,
                        tag_config=tag_config,
                        phrase_data=phrase_data,
                        month=month,
                        gemini_api_key=gemini_api_key,
                        slide=slide,
                        context=slide_context
                    )
                except Exception as e:
                    print(f"  ❌ 태그 처리 실패 ({tag_info['tag']}): {e}")
    
    # 저장
    prs.save(output_path)
    print(f"\n💾 리포트 저장 완료: {output_path}")