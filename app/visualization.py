"""
차트 생성 모듈
태그별로 필요한 차트를 생성하여 반환
"""
import pandas as pd
import numpy as np
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from app.config import Config

# 한글 폰트 설정
plt.rcParams['font.family'] = Config.CHART_FONT_FAMILY


def create_bubble_chart(
    df_phrases: pd.DataFrame,
    metric: str = '총 노출',
    top_n: int = 10,
    output_path: Optional[str] = None
) -> str:
    """
    버블 차트 생성
    
    Args:
        df_phrases: 프레이즈 데이터프레임
        metric: 사용할 지표 ('총 노출', '총 클릭', '평균 CTR')
        top_n: 상위 N개만 표시
        output_path: 출력 파일 경로 (None이면 자동 생성)
        
    Returns:
        생성된 차트 파일 경로
    """
    # 데이터 준비
    df_viz = df_phrases[df_phrases['프레이즈'] != '노이즈'].copy()
    df_viz = df_viz.sort_values(metric, ascending=False).head(top_n)
    
    # 버블 크기 정규화
    max_val = df_viz[metric].max()
    min_val = df_viz[metric].min()
    if max_val == min_val:
        bubble_sizes = np.full(len(df_viz), 50)
    else:
        bubble_sizes = 10 + (df_viz[metric] - min_val) / (max_val - min_val) * 90
    
    # 좌표 생성 (원형 배치)
    n = len(df_viz)
    angles = np.linspace(0, 2 * np.pi, n, endpoint=False)
    x_pos = np.cos(angles) * 3
    y_pos = np.sin(angles) * 3
    
    # 차트 생성
    fig, ax = plt.subplots(figsize=(10, 8))
    
    # CTR 범위 계산 (색상용)
    if '평균 CTR' in df_viz.columns:
        ctr_min = df_viz['평균 CTR'].min()
        ctr_max = df_viz['평균 CTR'].max()
        colors = plt.cm.RdYlGn((df_viz['평균 CTR'] - ctr_min) / (ctr_max - ctr_min))
    else:
        colors = plt.cm.viridis(np.linspace(0, 1, len(df_viz)))
    
    # 버블 그리기
    for i, (idx, row) in enumerate(df_viz.iterrows()):
        ax.scatter(
            x_pos[i], y_pos[i],
            s=bubble_sizes[i] * 50,
            c=[colors[i]],
            alpha=0.6,
            edgecolors='black',
            linewidths=1
        )
        # 텍스트 레이블
        ax.text(
            x_pos[i], y_pos[i],
            f"{row['프레이즈']}\n{row[metric]:,.0f}",
            ha='center', va='center',
            fontsize=9, weight='bold'
        )
    
    ax.set_title(f'{metric} 상위 {top_n}개 프레이즈', fontsize=16, weight='bold')
    ax.axis('off')
    plt.tight_layout()
    
    # 저장
    if output_path is None:
        output_path = f'chart_{metric}_{top_n}.png'
    
    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    
    plt.savefig(output_path, dpi=Config.CHART_DPI, bbox_inches='tight')
    plt.close()
    
    print(f"✅ 버블 차트 생성: {output_path}")
    return output_path


def find_chart_area_marker(slide, marker_text: str = "{{CHART_AREA}}") -> Optional[Tuple]:
    """
    슬라이드에서 차트 영역 마커를 찾아서 위치와 크기를 반환
    
    Parameters:
    -----------
    slide : pptx.slide.Slide
        검색할 슬라이드
    marker_text : str
        차트 영역을 표시하는 텍스트 (기본값: "{{CHART_AREA}}")
    
    Returns:
    --------
    tuple or None
        (left, top, width, height, marker_shape) in Inches, 또는 None (찾지 못한 경우)
    """
    for shape in slide.shapes:
        # 텍스트 프레임이 있는 shape 확인
        if hasattr(shape, 'text_frame'):
            if marker_text in shape.text:
                # 이 shape의 위치와 크기 반환
                return (shape.left, shape.top, shape.width, shape.height, shape)
        
        # shape 이름으로도 찾기 (PowerPoint에서 shape 이름 지정 가능)
        if hasattr(shape, 'name') and marker_text in shape.name:
            return (shape.left, shape.top, shape.width, shape.height, shape)
    
    return None


def insert_bubble_chart_shape_to_slide(
    slide,
    df_phrases: pd.DataFrame,
    metric: str = '총 노출',
    top_n: int = 10,
    chart_area: Optional[Tuple] = None,
    marker_text: str = "{{CHART_AREA}}",
    remove_marker: bool = True,
    color_metric: Optional[str] = None
) -> None:
    """
    슬라이드에 버블 차트를 원(shape)으로 직접 그리기
    
    Args:
        slide: pptx 슬라이드 객체
        df_phrases: 프레이즈 데이터프레임
        metric: 버블 크기 기준 지표 ('총 노출', '총 클릭', '평균 CTR')
        top_n: 상위 N개만 표시
        chart_area: 차트 영역 (left, top, width, height) - None이면 마커 찾기
        marker_text: 마커 텍스트
        remove_marker: 마커 shape 삭제 여부
        color_metric: 색상 기준 지표 (None이면 CTR 사용)
    """
    # 데이터 준비
    df_viz = df_phrases[df_phrases['프레이즈'] != '노이즈'].copy()
    df_viz = df_viz.sort_values(metric, ascending=False).head(top_n)
    
    if len(df_viz) == 0:
        print("⚠️ 표시할 데이터가 없습니다.")
        return
    
    # 1. 차트 영역 결정
    marker_shape = None
    if chart_area is None:
        # 마커 찾기 시도
        marker_result = find_chart_area_marker(slide, marker_text)
        if marker_result:
            chart_left, chart_top, chart_width, chart_height, marker_shape = marker_result
            print(f"✅ 차트 영역 마커 발견: {marker_text}")
        else:
            # 기본 위치 사용
            chart_left = Inches(0.5)
            chart_top = Inches(1.5)
            chart_width = Inches(9.0)
            chart_height = Inches(5.5)
            print(f"⚠️ 마커를 찾지 못했습니다. 기본 위치 사용")
    else:
        chart_left, chart_top, chart_width, chart_height = chart_area
    
    # 2. 마커 제거
    if remove_marker and marker_shape:
        slide.shapes._spTree.remove(marker_shape._element)
        print(f"✅ 마커 제거 완료")
    
    # 3. 버블 크기 정규화
    max_val = df_viz[metric].max()
    min_val = df_viz[metric].min()
    if max_val == min_val:
        bubble_sizes = np.full(len(df_viz), 50)
    else:
        bubble_sizes = 10 + (df_viz[metric] - min_val) / (max_val - min_val) * 90
    
    # 4. 색상 기준 지표 결정
    if color_metric is None:
        color_metric = '평균 CTR' if '평균 CTR' in df_viz.columns else metric
    
    color_min = df_viz[color_metric].min()
    color_max = df_viz[color_metric].max()
    
    # 5. 좌표 변환 함수
    def to_canvas_x(x01):
        return chart_left + chart_width * ((x01 + 2) / 10.0)
    
    def to_canvas_y(y01):
        # PowerPoint는 좌상단 원점(아래로 +)
        return chart_top + chart_height * (1.0 - ((y01 + 2) / 10.0))
    
    def size_to_diameter(size_val):
        d_min, d_max = 0.7, 1.8
        return Inches(d_min + (size_val - 10) / (100 - 10) * (d_max - d_min))
    
    def auto_text_color(rgb):
        r, g, b = rgb[0], rgb[1], rgb[2]
        luminance = 0.2126*r + 0.7152*g + 0.0722*b
        return RGBColor(255, 255, 255) if luminance < 140 else RGBColor(0, 0, 0)
    
    # 6. 색상 팔레트 (하늘색-보라색 그라데이션)
    def get_color_by_value(value, val_min, val_max):
        if val_max == val_min:
            ratio = 0.5
        else:
            ratio = (value - val_min) / (val_max - val_min)
        
        sky_r, sky_g, sky_b = 135, 206, 250  # 하늘색
        purple_r, purple_g, purple_b = 138, 43, 226  # 보라색
        
        r = int(sky_r + (purple_r - sky_r) * ratio)
        g = int(sky_g + (purple_g - sky_g) * ratio)
        b = int(sky_b + (purple_b - sky_b) * ratio)
        
        return (r, g, b)
    
    # 7. 위치 계산 (원형 배치)
    np.random.seed(42)
    n = len(df_viz)
    angles = np.linspace(0, 2 * np.pi, n, endpoint=False)
    radius_variation = np.random.uniform(0.4, 0.9, n)
    x_pos = np.cos(angles) * radius_variation
    y_pos = np.sin(angles) * radius_variation
    
    # 좌표 범위를 0~6으로 정규화
    coord_range = 6.0
    coord_center = 3.0
    if x_pos.max() != x_pos.min():
        x_pos = coord_center - coord_range/2 + (x_pos - x_pos.min()) / (x_pos.max() - x_pos.min()) * coord_range
    if y_pos.max() != y_pos.min():
        y_pos = coord_center - coord_range/2 + (y_pos - y_pos.min()) / (y_pos.max() - y_pos.min()) * coord_range
    
    # 8. 버블을 큰 것부터(뒤에 그리면 가려져서) → 작은 것부터 그려서 시각적으로 위로
    bubble_sizes_array = bubble_sizes.values if hasattr(bubble_sizes, 'values') else np.array(bubble_sizes)
    order = sorted(range(len(df_viz)), key=lambda i: bubble_sizes_array[i], reverse=True)
    
    # 9. 버블 그리기
    for i in order:
        # 좌표/크기
        cx = to_canvas_x(float(x_pos[i]))
        cy = to_canvas_y(float(y_pos[i]))
        d = size_to_diameter(float(bubble_sizes_array[i]))
        
        # 원(타원) 추가: 좌상단 기준이라 중심 보정
        left = cx - d/2
        top = cy - d/2
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,  # 원형
            left=left, top=top, width=d, height=d
        )
        
        # 색상
        color_val = df_viz.iloc[i][color_metric]
        rgb = get_color_by_value(color_val, color_min, color_max)
        
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        
        line = shape.line
        line.color.rgb = RGBColor(255, 255, 255)
        line.width = Pt(0.75)
        
        # 텍스트 구성
        phrase = df_viz.iloc[i]['프레이즈']
        if len(phrase) > 15:
            if '·' in phrase:
                parts = phrase.split('·', 1)
                phrase = parts[0] + '\n·' + parts[1]
            elif len(phrase) > 20:
                mid = len(phrase) // 2
                phrase = phrase[:mid] + '\n' + phrase[mid:]
        
        # 지표 값 표시
        if metric == '평균 CTR':
            value_text = f"{df_viz.iloc[i][metric]:.1f}%"
        else:
            value_text = f"{int(df_viz.iloc[i][metric]):,}"
        
        label_text = f"{phrase}\n{value_text}"
        
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = label_text
        
        # 원 지름에 비례한 폰트 크기
        pt = max(9, min(18, int(d.inches * 8)))
        r1.font.size = Pt(pt)
        r1.font.bold = False
        
        # 대비에 맞춰 글자색 자동
        r1.font.color.rgb = auto_text_color(rgb)
    
    print(f"✅ 버블 차트 생성 완료: {len(df_viz)}개 버블")


def create_chart_for_tag(
    tag_name: str,
    phrase_data: List[Dict[str, Any]],
    config: Dict[str, Any],
    slide=None,
    marker_text: Optional[str] = None
) -> Optional[str]:
    """
    태그에 맞는 차트 생성 및 삽입
    
    Args:
        tag_name: 태그 이름
        phrase_data: 프레이즈 데이터
        config: 태그 설정
        slide: 슬라이드 객체 (None이면 이미지 파일로 저장)
        marker_text: 마커 텍스트 (슬라이드에 직접 삽입할 경우)
        
    Returns:
        슬라이드에 삽입한 경우 None, 이미지 파일인 경우 파일 경로
    """
    # phrase_data를 DataFrame으로 변환
    df_phrases = pd.DataFrame(phrase_data)
    
    # 태그별 차트 타입 결정 (chart_kind 또는 chart_type 지원)
    chart_kind = config.get('chart_kind') or config.get('chart_type', 'bubble')
    metric = config.get('metric', '총 노출')
    top_n = config.get('top_n', 10)
    
    # 영어 metric을 한글로 변환 (필요시)
    metric_map = {
        'impressions': '총 노출',
        'clicks': '총 클릭',
        'ctr': '평균 CTR'
    }
    if metric in metric_map:
        metric = metric_map[metric]
    
    # 색상 기준 지표 결정
    color_metric = config.get('color_metric')
    if color_metric and color_metric in metric_map:
        color_metric = metric_map[color_metric]
    
    if chart_kind == 'bubble':
        if slide is not None:
            # 슬라이드에 직접 삽입
            if marker_text is None:
                marker_text = f"{{{{{tag_name}}}}}"
            
            insert_bubble_chart_shape_to_slide(
                slide=slide,
                df_phrases=df_phrases,
                metric=metric,
                top_n=top_n,
                marker_text=marker_text,
                color_metric=color_metric
            )
            return None
        else:
            # 이미지 파일로 저장
            output_path = config.get('output_path', f'chart_{tag_name}.png')
            return create_bubble_chart(
                df_phrases=df_phrases,
                metric=metric,
                top_n=top_n,
                output_path=output_path
            )
    else:
        raise ValueError(f"지원하지 않는 차트 타입: {chart_kind}")

